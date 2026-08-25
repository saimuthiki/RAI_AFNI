# -*- coding: utf-8 -*-
"""Checks capability-matrix table cells for likely text overflow (QA script
skips table cells since GraphicFrame has no has_text_frame)."""
import os
from pptx import Presentation
from pptx.util import Emu

PATH = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))),
                     "deliverables", "AFNI_Responsible_AI_Framework.pptx")
EMU = 914400
prs = Presentation(PATH)

issues = []
table_slides = 0

for si, slide in enumerate(prs.slides, start=1):
    for shp in slide.shapes:
        if not shp.has_table:
            continue
        table_slides += 1
        table = shp.table
        n_rows = len(table.rows)
        for ci, col in enumerate(table.columns):
            col_w_in = col.width / EMU
            for ri in range(n_rows):
                row_h_in = table.rows[ri].height / EMU
                cell = table.cell(ri, ci)
                for p in cell.text_frame.paragraphs:
                    text = "".join(r.text for r in p.runs)
                    if not text:
                        continue
                    size_pt = 10
                    for r in p.runs:
                        if r.font.size:
                            size_pt = r.font.size.pt
                            break
                    avg_char_w = (size_pt / 72) * 0.52
                    usable_w = max(col_w_in - 0.09, 0.1)
                    chars_per_line = max(1, int(usable_w / avg_char_w))
                    est_lines = max(1, -(-len(text) // chars_per_line))
                    n_paras = sum(1 for pp in cell.text_frame.paragraphs if "".join(rr.text for rr in pp.runs))
                    line_h = (size_pt / 72) * 1.25
                    est_h = est_lines * line_h
                    if est_h > row_h_in * 0.95 / max(n_paras, 1) + 0.02:
                        issues.append(
                            f"slide {si} table row {ri} col {ci}: col_w={col_w_in:.2f}in row_h={row_h_in:.2f}in "
                            f"{size_pt}pt '{text[:40].encode('ascii','replace').decode()}' "
                            f"est_lines={est_lines} chars_per_line={chars_per_line}"
                        )

print(f"Tables found: {table_slides}")
print(f"Potential overflow cells: {len(issues)}")
for i in issues[:60]:
    print(" -", i)
