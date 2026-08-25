# -*- coding: utf-8 -*-
"""
AFNI Responsible AI Framework - PPTX builder.
Reads data/RAI_Repo_Reports.json and data/RAI_Synthesis.json (relative to the
project root) and generates AFNI_Responsible_AI_Framework.pptx at the root.
"""
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
_DATA_DIR = os.path.join(_ROOT, "data")

REPORTS_PATH = os.path.join(_DATA_DIR, "RAI_Repo_Reports.json")
SYNTHESIS_PATH = os.path.join(_DATA_DIR, "RAI_Synthesis.json")
# The deck is a client deliverable, so it lands in deliverables/ at the repo
# root rather than beside the pipeline that builds it.
_DELIVERABLES = os.path.join(os.path.dirname(_ROOT), "deliverables")
OUT_PATH = os.path.join(_DELIVERABLES, "AFNI_Responsible_AI_Framework.pptx")

# ---------------------------------------------------------------- THEME ----
NAVY = RGBColor(0x10, 0x24, 0x3E)
NAVY_DARK = RGBColor(0x0A, 0x17, 0x2B)
TEAL = RGBColor(0x1B, 0x9A, 0xAA)
BG_LIGHT = RGBColor(0xF6, 0xF8, 0xFB)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_MUTED = RGBColor(0x5B, 0x6B, 0x7D)
TEXT_SOFT_ON_NAVY = RGBColor(0xCF, 0xDA, 0xE6)
AMBER = RGBColor(0xE8, 0x8C, 0x00)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
RED_SOFT = RGBColor(0xC0, 0x39, 0x2B)
LINE_GREY = RGBColor(0xDD, 0xE3, 0xEA)

TENET_COLORS = {
    "Privacy": RGBColor(0x2E, 0x86, 0xAB),
    "Security": RGBColor(0xD6, 0x28, 0x28),
    "Fairness & Bias": RGBColor(0x6A, 0x4C, 0x93),
    "Explainability & Transparency": RGBColor(0xE0, 0x8E, 0x45),
    "Profanity / Content Safety": RGBColor(0xE7, 0x6F, 0x51),
    "Hallucination / Reliability": RGBColor(0x21, 0x8A, 0x7E),
    "Accountability": RGBColor(0x4A, 0x63, 0x80),
}
TENET_ORDER = list(TENET_COLORS.keys())

ROLE_COLORS = {
    "Guardrail Development": RGBColor(0x1B, 0x9A, 0xAA),
    "Vulnerability / Red-Team Testing": RGBColor(0xC0, 0x39, 0x2B),
    "Both": RGBColor(0x6A, 0x4C, 0x93),
}

FONT_HEAD = "Calibri"
FONT_BODY = "Calibri"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5


# ------------------------------------------------------------- HELPERS -----
def new_prs():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    return prs


def blank_slide(prs, bg=BG_LIGHT):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    rect.shadow.inherit = False
    # send to back
    spTree = slide.shapes._spTree
    spTree.remove(rect._element)
    spTree.insert(2, rect._element)
    return slide


def add_rect(slide, x, y, w, h, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = LINE_GREY
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_rounded(slide, x, y, w, h, color, radius=0.06, line=False, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = line_color or LINE_GREY
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=14, color=TEXT_DARK, bold=False, italic=False,
             align=PP_ALIGN.LEFT, font=FONT_BODY, anchor=MSO_ANCHOR.TOP, line_spacing=1.0,
             wrap=True, shrink=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
    return box


def add_bullets(slide, x, y, w, h, items, size=13, color=TEXT_DARK, font=FONT_BODY,
                 bullet_color=None, space_after=6, bold_first=False, line_spacing=1.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    bullet_color = bullet_color or TEAL
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        r1 = p.add_run()
        r1.text = "\u25AA  "
        r1.font.size = Pt(size)
        r1.font.color.rgb = bullet_color
        r1.font.bold = True
        r1.font.name = font
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
        r2.font.name = font
    return box


def add_footer(slide, page_no, section=""):
    add_rect(slide, 0, SLIDE_H_IN - 0.32, SLIDE_W_IN, 0.32, NAVY)
    add_text(slide, 0.4, SLIDE_H_IN - 0.32, 6, 0.32, "AFNI Responsible AI Framework", size=9,
              color=TEXT_SOFT_ON_NAVY, anchor=MSO_ANCHOR.MIDDLE, font=FONT_BODY)
    if section:
        add_text(slide, 6.4, SLIDE_H_IN - 0.32, 4.0, 0.32, section, size=9,
                  color=TEXT_SOFT_ON_NAVY, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, font=FONT_BODY)
    add_text(slide, SLIDE_W_IN - 1.0, SLIDE_H_IN - 0.32, 0.6, 0.32, str(page_no), size=9,
              color=TEXT_SOFT_ON_NAVY, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT, font=FONT_BODY)


def add_header(slide, kicker, title, accent=TEAL):
    add_rect(slide, 0, 0, 0.14, SLIDE_H_IN - 0.32, accent)
    add_text(slide, 0.55, 0.32, 10, 0.3, kicker.upper(), size=12, color=accent, bold=True, font=FONT_HEAD)
    add_text(slide, 0.55, 0.62, 12.2, 0.7, title, size=26, color=NAVY, bold=True, font=FONT_HEAD)
    add_rect(slide, 0.55, 1.32, 2.2, 0.035, accent)


def pill(slide, x, y, text, color, w=None, size=10.5, text_color=WHITE):
    w = w if w else (0.13 + 0.093 * len(text) + 0.28)
    shp = add_rounded(slide, x, y, w, 0.28, color, radius=0.5)
    add_text(slide, x, y, w, 0.28, text, size=size, color=text_color, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False, font=FONT_HEAD)
    return w


def tenet_pills(slide, x, y, tenets, max_w=12.2, size=10):
    cx = x
    cy = y
    row_h = 0.34
    for t in tenets:
        color = TENET_COLORS.get(t, TEAL)
        short = t.split(" / ")[0].split(" & ")[0].split(" (")[0]
        w = 0.16 + 0.078 * len(short) + 0.2
        if cx + w > x + max_w:
            cx = x
            cy += row_h
        pill(slide, cx, cy, short, color, w=w, size=size)
        cx += w + 0.12
    return cy + row_h


def badge(slide, x, y, text, color, w=1.7, h=0.3, size=10.5):
    shp = add_rounded(slide, x, y, w, h, color, radius=0.22)
    add_text(slide, x, y, w, h, text, size=size, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE, wrap=False, font=FONT_HEAD)


def add_table(slide, x, y, w, h, headers, rows, col_widths=None, header_bg=NAVY,
              header_fg=WHITE, font_size=11, row_h=None, header_h=0.4, zebra=True,
              body_fg=TEXT_DARK):
    n_cols = len(headers)
    n_rows = len(rows) + 1
    graphic_frame = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = graphic_frame.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(w * cw / total)
    # header
    for j, htext in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = htext
        r.font.size = Pt(font_size + 0.5)
        r.font.bold = True
        r.font.color.rgb = header_fg
        r.font.name = FONT_HEAD
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (RGBColor(0xF0, 0xF4, 0xF8) if (zebra and i % 2 == 1) else WHITE)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(font_size)
            r.font.color.rgb = body_fg
            r.font.name = FONT_BODY
    # row heights
    table.rows[0].height = Inches(header_h)
    if row_h:
        for i in range(1, n_rows):
            table.rows[i].height = Inches(row_h)
    return graphic_frame


def set_cell_text_color_by_value(cell, value_map):
    pass  # placeholder for future conditional coloring


# ------------------------------------------------------- FLOWCHART HELPERS --
def add_line(slide, x1, y1, x2, y2, color=TEXT_MUTED, width=1.5, dashed=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    if dashed:
        d = conn.line._get_or_add_ln()
        dash = d.makeelement(qn('a:prstDash'), {'val': 'dash'})
        d.append(dash)
    conn.shadow.inherit = False
    return conn


def add_arrowhead(slide, x, y, direction, color=TEXT_MUTED, size=0.09):
    shp = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x - size / 2), Inches(y - size / 2),
                                  Inches(size), Inches(size))
    rot = {"right": 90, "down": 180, "left": 270, "up": 0}[direction]
    shp.rotation = rot
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def flow_arrow(slide, x1, y1, x2, y2, color=TEXT_MUTED, width=1.5, dashed=False, label=None, label_size=8.5):
    """Draws a straight connector with an arrowhead at (x2,y2). Assumes axis-aligned (h or v)."""
    add_line(slide, x1, y1, x2, y2, color=color, width=width, dashed=dashed)
    if abs(x2 - x1) >= abs(y2 - y1):
        direction = "right" if x2 >= x1 else "left"
    else:
        direction = "down" if y2 >= y1 else "up"
    add_arrowhead(slide, x2, y2, direction, color=color)
    if label:
        mx, my = (x1 + x2) / 2, (y1 + y2) / 2
        add_text(slide, mx - 0.6, my - 0.16, 1.2, 0.18, label, size=label_size, color=color,
                  align=PP_ALIGN.CENTER, font=FONT_BODY, wrap=False)


def flow_box(slide, x, y, w, h, text, subtext=None, fill=CARD_BG, line_color=LINE_GREY, text_color=TEXT_DARK,
             size=10.5, sub_size=8.3, shape=MSO_SHAPE.ROUNDED_RECTANGLE, bold=True):
    shp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            shp.adjustments[0] = 0.12
        except Exception:
            pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line_color
    shp.line.width = Pt(1)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = text_color
    r.font.name = FONT_HEAD
    if subtext:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = subtext
        r2.font.size = Pt(sub_size)
        r2.font.bold = False
        r2.font.color.rgb = TEXT_MUTED
        r2.font.name = FONT_BODY
    return shp


if __name__ == "__main__":
    print("This module provides helpers only; run build_deck.py to generate the PPTX.")
