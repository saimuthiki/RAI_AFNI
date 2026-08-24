# -*- coding: utf-8 -*-
"""Heuristic QA: flags text frames that likely overflow their box, and slides
whose shapes fall outside the slide bounds. No rendering engine available,
so this is an approximation (avg char width ~0.5x font size in points)."""
import os
from pptx import Presentation
from pptx.util import Emu

PATH = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                     "AFNI_Responsible_AI_Framework.pptx")
EMU_PER_IN = 914400

prs = Presentation(PATH)
slide_w = prs.slide_width / EMU_PER_IN
slide_h = prs.slide_height / EMU_PER_IN

issues = []

for si, slide in enumerate(prs.slides, start=1):
    for shp in slide.shapes:
        try:
            x, y, w, h = shp.left / EMU_PER_IN, shp.top / EMU_PER_IN, shp.width / EMU_PER_IN, shp.height / EMU_PER_IN
        except TypeError:
            continue
        if x < -0.05 or y < -0.05 or (x + w) > slide_w + 0.05 or (y + h) > slide_h + 0.05:
            issues.append(f"slide {si}: shape out of bounds x={x:.2f} y={y:.2f} w={w:.2f} h={h:.2f}")

        if not shp.has_text_frame:
            continue
        tf = shp.text_frame
        total_chars = sum(len(p.text) for p in tf.paragraphs)
        if total_chars == 0:
            continue
        # estimate line height in inches from first run's font size (pt), fallback 12pt
        size_pt = 12
        for p in tf.paragraphs:
            for r in p.runs:
                if r.font.size:
                    size_pt = r.font.size.pt
                    break
            break
        line_h_in = (size_pt / 72) * 1.35
        avg_char_w_in = (size_pt / 72) * 0.52
        usable_w = max(w, 0.1)
        chars_per_line = max(1, int(usable_w / avg_char_w_in))
        n_paragraphs = len(tf.paragraphs)
        est_lines = 0
        for p in tf.paragraphs:
            t = p.text
            est_lines += max(1, -(-len(t) // chars_per_line)) if t else 1
        est_height = est_lines * line_h_in
        if est_height > h + 0.12:
            preview = " / ".join(p.text[:40] for p in tf.paragraphs if p.text)[:90]
            issues.append(
                f"slide {si}: possible overflow - box h={h:.2f}in est_text_h={est_height:.2f}in "
                f"({size_pt:.0f}pt, {total_chars} chars) :: {preview}"
            )

print(f"Slides: {len(prs.slides._sldIdLst)}")
print(f"Issues found: {len(issues)}")
for i in issues:
    print(" -", i.encode("ascii", "replace").decode("ascii"))
