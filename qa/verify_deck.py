# -*- coding: utf-8 -*-
import os
import sys

sys.path.insert(0, os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "helpers"))

from pptx import Presentation
from repo_slide_content import REPO_SLIDES

PATH = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                     "AFNI_Responsible_AI_Framework.pptx")
prs = Presentation(PATH)

all_text = []
for i, slide in enumerate(prs.slides, start=1):
    texts = []
    for shp in slide.shapes:
        if shp.has_text_frame:
            t = shp.text_frame.text.strip()
            if t:
                texts.append(t)
    all_text.append((i, texts))

full_doc_text = "\n".join(t for _, texts in all_text for t in texts)

print(f"Total slides: {len(all_text)}\n")
print("=== Slide-by-slide first line (title-ish) ===")
for i, texts in all_text:
    title_guess = texts[1] if len(texts) > 1 else (texts[0] if texts else "(empty)")
    print(f"{i:3d}: {title_guess[:80]}")

print("\n=== Coverage check: all 23 repo display names present? ===")
missing_repos = []
for r in REPO_SLIDES:
    if r["display_name"] not in full_doc_text:
        missing_repos.append(r["display_name"])
print("Missing:", missing_repos if missing_repos else "NONE - all 23 present")

print("\n=== Coverage check: all 7 tenets present? ===")
TENETS = ["Privacy", "Security", "Fairness & Bias", "Explainability & Transparency",
          "Profanity / Content Safety", "Hallucination / Reliability", "Accountability"]
missing_tenets = [t for t in TENETS if t not in full_doc_text]
print("Missing:", missing_tenets if missing_tenets else "NONE - all 7 present")
